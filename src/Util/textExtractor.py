import os
import tempfile
from typing import List, Union
from tika import parser
from pdf2image import convert_from_path
from PIL import Image
from pptx import Presentation
from PyPDF2 import PdfReader
import streamlit as st
from streamlit.runtime.uploaded_file_manager import UploadedFile


class DocumentTextExtractor:
    def extract_text(self, file: Union[str, UploadedFile]) -> str:
        """
        Extract text from a given file (PPT, image, or PDF).

        :param file: Path to the input file or UploadedFile object
        :return: Extracted text as a string
        """
        if isinstance(file, UploadedFile):
            with tempfile.NamedTemporaryFile(
                delete=False, suffix=file.name
            ) as tmp_file:
                tmp_file.write(file.getvalue())
                file_path = tmp_file.name
        else:
            file_path = file

        file_extension = os.path.splitext(file_path)[1].lower()

        if file_extension in [".ppt", ".pptx"]:
            text = self._extract_text_from_ppt(file_path)
        elif file_extension in [".jpg", ".jpeg", ".png", ".bmp", ".gif"]:
            text = self._extract_text_from_image(file_path)
        elif file_extension == ".pdf":
            text = self._extract_text_from_pdf(file_path)
        else:
            raise ValueError(f"Unsupported file format: {file_extension}")

        # Clean up temporary file if created
        if isinstance(file, UploadedFile):
            os.unlink(file_path)

        return text

    def _extract_text_from_ppt(self, ppt_path: str) -> str:
        """
        Extract text from PowerPoint file.
        """
        prs = Presentation(ppt_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    def _extract_text_from_image(self, image_path: str) -> str:
        """
        Extract text from image using Tika.
        """
        raw = parser.from_file(image_path)
        return raw["content"] or "No text extracted from image."

    def _extract_text_from_pdf(self, pdf_path: str) -> str:
        """
        Extract text from PDF using PyPDF2.
        """
        with open(pdf_path, "rb") as file:
            reader = PdfReader(file)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
        return text

    def extract_text_from_multiple_files(
        self, files: List[Union[str, UploadedFile]]
    ) -> str:
        """
        Extract text from multiple files and combine the results.

        :param files: List of paths to input files or UploadedFile objects
        :return: Combined extracted text as a string, separated by filenames and underscores
        """
        combined_text = ""
        for file in files:
            if isinstance(file, UploadedFile):
                filename = file.name
            else:
                filename = os.path.basename(file)

            try:
                text = self.extract_text(file)
                combined_text += f"\n\n{'_' * 40}\n{filename}\n{'_' * 40}\n\n{text}\n"
            except Exception as e:
                combined_text += f"\n\n{'_' * 40}\n{filename}\n{'_' * 40}\n\nError processing file: {str(e)}\n"

        return combined_text.strip()
